# services/input_parser.py

from pathlib import Path

from services.document_parser import extract_text as extract_document_text
from services.transcription import transcribe_audio, transcribe_image


# Supported file extensions
DOCUMENT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".txt",
    ".rtf",
    ".odt",
}

SPREADSHEET_EXTENSIONS = {
    ".xlsx",
    ".xls",
    ".csv",
}

PRESENTATION_EXTENSIONS = {
    ".pptx",
}

IMAGE_EXTENSIONS = {
    ".jpg",
    ".jpeg",
    ".png",
    ".webp",
    ".bmp",
    ".tiff",
    ".tif",
}

AUDIO_EXTENSIONS = {
    ".wav",
    ".mp3",
    ".m4a",
    ".flac",
    ".ogg",
    ".aac",
    ".webm",
}


def get_file_type(file_path):
    """
    Determine the type of an uploaded file.

    Returns:
        "document"
        "spreadsheet"
        "presentation"
        "image"
        "audio"
        "unsupported"
    """

    extension = Path(file_path).suffix.lower()

    if extension in DOCUMENT_EXTENSIONS:
        return "document"

    if extension in SPREADSHEET_EXTENSIONS:
        return "spreadsheet"

    if extension in PRESENTATION_EXTENSIONS:
        return "presentation"

    if extension in IMAGE_EXTENSIONS:
        return "image"

    if extension in AUDIO_EXTENSIONS:
        return "audio"

    return "unsupported"


def extract_input(file_path):
    """
    Extract useful text from an uploaded file.

    All supported input types are converted into text so that
    the resulting text can be passed into the SAATHI AI pipeline.

    Args:
        file_path: Path to uploaded file.

    Returns:
        dict containing:
            type
            text
            filename
    """

    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(
            f"File not found: {file_path}"
        )

    file_type = get_file_type(file_path)

    if file_type == "unsupported":
        raise ValueError(
            f"Unsupported file type: {path.suffix.lower()}"
        )

    # Documents
    if file_type == "document":
        text = extract_document_text(file_path)

    # Audio
    elif file_type == "audio":
        text = transcribe_audio(file_path)

    # These will be implemented by dedicated parsers
    elif file_type == "spreadsheet":
        text = extract_spreadsheet_text(file_path)

    elif file_type == "presentation":
        text = extract_presentation_text(file_path)

    elif file_type == "image":
        text = extract_image_text(file_path)

    else:
        raise ValueError(
            f"Unsupported input type: {file_type}"
        )

    if not text or not text.strip():
        raise ValueError(
            "No readable text could be extracted from the file."
        )

    return {
        "type": file_type,
        "filename": path.name,
        "text": text.strip()
    }


def extract_spreadsheet_text(file_path):
    """
    Extract spreadsheet contents as text.

    Supports:
        XLSX
        XLS
        CSV
    """

    import pandas as pd

    path = Path(file_path)
    extension = path.suffix.lower()

    if extension == ".csv":
        dataframe = pd.read_csv(file_path)

    elif extension in {".xlsx", ".xls"}:
        dataframe = pd.read_excel(file_path)

    else:
        raise ValueError(
            f"Unsupported spreadsheet type: {extension}"
        )

    return dataframe.to_string(index=False)


def extract_presentation_text(file_path):
    """
    Extract text from PowerPoint presentations.
    """

    from pptx import Presentation

    presentation = Presentation(file_path)

    text = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            text.append(
                f"Slide {slide_number}:\n"
                + "\n".join(slide_text)
            )

    return "\n\n".join(text)


def extract_image_text(file_path):
    """
    Extract text from an image using handwriting-aware OCR preprocessing.
    """
    return transcribe_image(file_path)
